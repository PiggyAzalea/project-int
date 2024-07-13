"""Module docstring for app.py"""

import os
import mimetypes
import shutil
from flask import Flask, render_template, request, redirect, url_for, flash, send_file, jsonify
from werkzeug.utils import secure_filename
import mammoth
from pptx import Presentation
import fitz  # PyMuPDF
from jinja2 import Environment, FileSystemLoader, select_autoescape, Undefined
import nbformat
from jupyter_client import KernelManager
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import contextlib

app = Flask(__name__, static_url_path='/static', static_folder='static')
upload_folder = os.path.abspath('templates/uploads')
app.config['UPLOAD_FOLDER'] = upload_folder
app.secret_key = 'supersecretkey'
file_index = {}

# Global dictionary to store variables
global_vars = {}

def extract_text_and_images_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    images = []
    image_dir = os.path.splitext(file_path)[0] + "_images"

    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_name = f"image_{len(images) + 1}.png"
                image_path = os.path.join(image_dir, image_name)
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                images.append({
                    'image': image,
                    'image_name': image_name,
                    'image_path': image_path
                })
            if hasattr(shape, "text"):
                text.append(shape.text)

    return {"text": "\n".join(text), "images": images}

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)
    file = request.files['file']
    if file.filename == '':
        flash('No selected file')
        return redirect(request.url)
    if file:
        filename = secure_filename(file.filename)
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
        file_index[filename] = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        return redirect(url_for('uploaded_file', filename=filename))
    return redirect(request.url)

@app.route('/uploads/<filename>')
def uploaded_file(filename):
    if filename not in file_index:
        flash('File not found')
        return redirect(url_for('index'))
    file_path = file_index[filename]
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        return send_file(file_path, mimetype=mime_type)
    else:
        flash('Unsupported file type')
        return redirect(url_for('index'))

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'filename' not in request.form:
        flash('No filename provided')
        return redirect(url_for('index'))
    filename = request.form['filename']
    if filename not in file_index:
        flash('File not found')
        return redirect(url_for('index'))
    file_path = file_index[filename]
    converted_data = extract_text_and_images_from_pptx(file_path)
    return jsonify(converted_data)

@app.route('/execute', methods=['POST'])
def run_python():
    global execution_environment
    code = request.form['code']
    output = io.StringIO()
    error = io.StringIO()
    try:
        with contextlib.redirect_stdout(output), contextlib.redirect_stderr(error):
            exec(code, globals())
    except Exception as e:
        error.write(str(e))
    return jsonify({
        'output': output.getvalue(),
        'error': error.getvalue()
    })

if __name__ == '__main__':
    app.run(debug=True)
